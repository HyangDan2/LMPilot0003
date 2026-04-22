from __future__ import annotations

from pathlib import Path

from src.models.schemas import PresentationPlan
from src.utils.io import ensure_dir
from src.utils.paths import slugify


class RendererError(Exception):
    """Raised when a presentation cannot be rendered."""


class PptxRenderer:
    """Deterministic PowerPoint renderer for validated presentation plans."""

    def render(self, plan: PresentationPlan, output_dir: Path, output_filename: str | None = None) -> Path:
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
            from pptx.util import Inches, Pt  # type: ignore[import-not-found]
        except Exception as exc:
            raise RendererError("PPTX rendering requires python-pptx.") from exc

        output_dir = ensure_dir(output_dir)
        presentation = Presentation()

        title_slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = plan.title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"Audience: {plan.target_audience}"

        content_layout = presentation.slide_layouts[5]
        for slide_plan in plan.slides:
            slide = presentation.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_plan.slide_title
            body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(4.8))
            frame = body.text_frame
            frame.clear()
            paragraph = frame.paragraphs[0]
            paragraph.text = slide_plan.purpose
            paragraph.font.size = Pt(20)

            if slide_plan.source_refs:
                refs = frame.add_paragraph()
                refs.text = "Sources: " + ", ".join(slide_plan.source_refs)
                refs.font.size = Pt(12)
            if slide_plan.image_refs:
                refs = frame.add_paragraph()
                refs.text = "Image refs: " + ", ".join(slide_plan.image_refs)
                refs.font.size = Pt(12)

        filename = output_filename or f"{slugify(plan.title, 'presentation')}.pptx"
        output_path = output_dir / filename
        try:
            presentation.save(str(output_path))
        except Exception as exc:
            raise RendererError(f"Failed to save PPTX: {exc}") from exc
        return output_path

