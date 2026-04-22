from __future__ import annotations

import os
import shutil
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Literal


ImageMode = Literal["auto", "ocr", "caption", "ocr_cv"]

SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".json",
    ".yaml",
    ".yml",
    ".csv",
    ".py",
    ".log",
    ".pdf",
    ".docx",
    ".pptx",
    ".png",
    ".jpg",
    ".jpeg",
    ".bmp",
    ".webp",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".webp"}
PLAIN_TEXT_EXTENSIONS = SUPPORTED_EXTENSIONS - {".pdf", ".docx", ".pptx"} - IMAGE_EXTENSIONS


class AttachmentError(Exception):
    pass


@dataclass(frozen=True)
class ExtractedAttachment:
    filename: str
    path: str
    file_type: str
    extracted_text: str


def list_supported_files_in_folder(folder_path: str) -> list[Path]:
    root = Path(folder_path).expanduser().resolve()
    if not root.is_dir():
        raise AttachmentError(f"Folder not found: {root}")

    paths: list[Path] = []
    for file_path in sorted(root.iterdir()):
        if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
            paths.append(file_path.resolve())

    if not paths:
        raise AttachmentError(f"No supported files found in folder: {root}")
    return paths


def validate_attachment_path(path: str) -> Path:
    file_path = Path(path).expanduser().resolve()
    suffix = file_path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise AttachmentError(f"Unsupported file type: {suffix or '(none)'}")
    if not file_path.is_file():
        raise AttachmentError(f"File not found: {file_path}")
    return file_path


def extract_text_from_file(path: str, image_mode: ImageMode = "auto") -> ExtractedAttachment:
    file_path = validate_attachment_path(path)
    suffix = file_path.suffix.lower()

    if suffix in PLAIN_TEXT_EXTENSIONS:
        text = _read_text_file(file_path)
    elif suffix == ".pdf":
        text = _read_pdf(file_path)
    elif suffix == ".docx":
        text = _read_docx(file_path)
    elif suffix == ".pptx":
        text = _read_pptx(file_path)
    else:
        text = _read_image(file_path, image_mode)

    if not text.strip():
        raise AttachmentError(f"No extractable text found in {file_path.name}")
    return ExtractedAttachment(file_path.name, str(file_path), suffix.lstrip("."), text.strip())


def format_attachment_context(attachments: list[dict[str, str]]) -> str:
    blocks: list[str] = []
    for attachment in attachments:
        filename = str(attachment.get("filename", "attached file"))
        file_type = str(attachment.get("file_type", "unknown"))
        text = str(attachment.get("extracted_text", "")).strip()
        if text:
            blocks.append(
                f"File name: {filename}\n"
                f"File type: {file_type}\n\n"
                f"Extracted content:\n{text}"
            )
    return "\n\n".join(blocks)


def format_user_text_with_attachments(user_text: str, attachment_context: str) -> str:
    if not attachment_context.strip():
        return user_text
    return f"{attachment_context.strip()}\n\nUser message:\n{user_text}"


def _read_text_file(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp949", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_bytes().decode("utf-8", errors="replace")


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore[import-not-found]
    except Exception as exc:
        raise AttachmentError("PDF extraction requires the optional 'pypdf' package.") from exc

    try:
        reader = PdfReader(str(path))
        pages: list[str] = []
        for index, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text() or ""
            pages.append(f"[Page {index}]\n{page_text.strip()}")
        return "\n\n".join(page for page in pages if page.strip())
    except Exception as exc:
        raise AttachmentError(f"Failed to extract PDF text: {exc}") from exc


def _read_docx(path: Path) -> str:
    try:
        from docx import Document  # type: ignore[import-not-found]
    except Exception as exc:
        raise AttachmentError("DOCX extraction requires the optional 'python-docx' package.") from exc

    try:
        document = Document(str(path))
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        return "\n".join(paragraphs)
    except Exception as exc:
        raise AttachmentError(f"Failed to extract DOCX text: {exc}") from exc


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except Exception as exc:
        raise AttachmentError("PPTX extraction requires the optional 'python-pptx' package.") from exc

    try:
        presentation = Presentation(str(path))
        slides: list[str] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = paragraph.text.strip()
                    if paragraph_text:
                        texts.append(paragraph_text)
            if texts:
                slides.append(f"[Slide {slide_number}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as exc:
        raise AttachmentError(f"Failed to extract PPTX text: {exc}") from exc


def _read_image(path: Path, image_mode: ImageMode) -> str:
    try:
        from PIL import Image, ImageFilter, ImageStat  # type: ignore[import-not-found]
    except Exception as exc:
        raise AttachmentError("Image preprocessing requires the optional 'Pillow' package.") from exc

    try:
        with Image.open(path) as image:
            image.load()
            metadata = _image_metadata(image, ImageStat, ImageFilter)
            normalized_mode = _normalize_image_mode(image_mode)
            parts = [f"Image preprocessing mode: {normalized_mode}"]

            if normalized_mode in {"auto", "ocr", "ocr_cv"}:
                parts.append(_ocr_image(image, path))
            if normalized_mode in {"auto", "caption"}:
                parts.append(_caption_image(image, metadata))
            if normalized_mode in {"auto", "ocr_cv"}:
                parts.append(_format_metadata(metadata))

            return "\n\n".join(part for part in parts if part.strip())
    except AttachmentError:
        raise
    except Exception as exc:
        raise AttachmentError(f"Failed to preprocess image: {exc}") from exc


def _normalize_image_mode(image_mode: str) -> ImageMode:
    normalized = image_mode.strip().lower().replace(" ", "_").replace("+", "_")
    if normalized in {"ocr_only", "ocr"}:
        return "ocr"
    if normalized in {"caption_only", "caption"}:
        return "caption"
    if normalized in {"ocr_cv", "ocr_cv_analysis", "ocr_analysis"}:
        return "ocr_cv"
    return "auto"


def _ocr_image(image, path: Path) -> str:
    tesseract_cmd = os.environ.get("TESSERACT_CMD", "").strip()
    try:
        import pytesseract  # type: ignore[import-not-found]
    except Exception:
        cli_text = _ocr_with_tesseract_cli(path, tesseract_cmd)
        if cli_text is not None:
            return cli_text
        return (
            "OCR text:\n"
            "OCR unavailable: install the Python package 'pytesseract' and the native Tesseract OCR engine. "
            "If Tesseract is not on PATH, set TESSERACT_CMD to the full executable path."
        )

    try:
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
        text = pytesseract.image_to_string(image).strip()
    except Exception as exc:
        cli_text = _ocr_with_tesseract_cli(path, tesseract_cmd)
        if cli_text is not None:
            return cli_text
        return (
            "OCR text:\n"
            f"OCR failed: {exc}. Confirm the native Tesseract OCR engine is installed and visible on PATH, "
            "or set TESSERACT_CMD to the executable path."
        )
    return f"OCR text:\n{text or '(no readable OCR text detected)'}"


def _ocr_with_tesseract_cli(path: Path, tesseract_cmd: str = "") -> str | None:
    command = tesseract_cmd or shutil.which("tesseract")
    if not command:
        return None
    try:
        completed = subprocess.run(
            [command, str(path), "stdout"],
            check=False,
            capture_output=True,
            text=True,
            timeout=60,
        )
    except Exception:
        return None
    if completed.returncode != 0:
        return None
    text = completed.stdout.strip()
    return f"OCR text:\n{text or '(no readable OCR text detected)'}"


def _caption_image(image, metadata: dict[str, object]) -> str:
    width = int(metadata["width"])
    height = int(metadata["height"])
    brightness = float(metadata["brightness"])
    contrast = float(metadata["contrast"])
    orientation = "landscape" if width > height else "portrait" if height > width else "square"
    tone = "bright" if brightness >= 180 else "dark" if brightness <= 70 else "moderately lit"
    detail = "high contrast" if contrast >= 64 else "low contrast" if contrast <= 24 else "moderate contrast"
    return (
        "Image caption:\n"
        f"A {orientation} {width}x{height} image that appears {tone} with {detail}. "
        "No semantic vision-captioning model is configured, so this caption is based on local image statistics."
    )


def _image_metadata(image, image_stat, image_filter) -> dict[str, object]:
    grayscale = image.convert("L")
    stats = image_stat.Stat(grayscale)
    edges = grayscale.filter(image_filter.FIND_EDGES)
    edge_stats = image_stat.Stat(edges)
    width, height = image.size
    return {
        "width": width,
        "height": height,
        "mode": image.mode,
        "format": image.format or "unknown",
        "brightness": round(float(stats.mean[0]), 2),
        "contrast": round(float(stats.stddev[0]), 2),
        "edge_density": round(float(edge_stats.mean[0]), 2),
    }


def _format_metadata(metadata: dict[str, object]) -> str:
    lines = ["Image metadata:"]
    for key, value in metadata.items():
        lines.append(f"- {key}: {value}")
    return "\n".join(lines)
