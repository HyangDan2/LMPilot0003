from __future__ import annotations

from pathlib import Path
from typing import Any

from src.ingestion.parsers.base import DocumentParser, ParserError
from src.models.schemas import Asset, ParsedDocument, Section


class PptxParser(DocumentParser):
    """MVP PPTX parser that extracts slide text and image-like shape metadata."""

    file_type = "pptx"

    def parse(self, path: Path) -> ParsedDocument:
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-not-found]
        except Exception as exc:
            raise ParserError("PPTX parsing requires python-pptx.") from exc

        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise ParserError(f"Failed to open PPTX {path}: {exc}") from exc

        doc_id = self.doc_id(path)
        sections: list[Section] = []
        assets: list[Asset] = []
        slide_texts: list[str] = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts = [_shape_text(shape) for shape in slide.shapes]
            texts = [text for text in texts if text]
            title = _slide_title(slide, slide_number)
            slide_text = "\n".join(texts).strip()
            slide_texts.append(slide_text)

            section_assets: list[Asset] = []
            for shape_index, shape in enumerate(slide.shapes, start=1):
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    asset = Asset(
                        asset_id=f"{doc_id}-slide-{slide_number}-asset-{shape_index}",
                        kind="image",
                        source_file=str(path.resolve()),
                        page_or_slide=slide_number,
                        path="",
                        caption=getattr(shape, "name", ""),
                        metadata={
                            "shape_name": getattr(shape, "name", ""),
                            "width": int(getattr(shape, "width", 0)),
                            "height": int(getattr(shape, "height", 0)),
                        },
                    )
                    section_assets.append(asset)
                    assets.append(asset)

            sections.append(
                Section(
                    section_id=f"{doc_id}-slide-{slide_number}",
                    title=title,
                    level=1,
                    text=slide_text,
                    page_or_slide=slide_number,
                    assets=section_assets,
                    metadata={"parser": "pptx", "shape_count": len(slide.shapes)},
                )
            )

        return ParsedDocument(
            doc_id=doc_id,
            file_path=str(path.resolve()),
            file_type=self.file_type,
            title=self.title_from_path(path),
            text="\n\n".join(text for text in slide_texts if text).strip(),
            sections=sections,
            assets=assets,
            metadata={"slide_count": len(presentation.slides)},
        )


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())


def _slide_title(slide: Any, slide_number: int) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None:
        title = _shape_text(title_shape)
        if title:
            return title.splitlines()[0]
    return f"Slide {slide_number}"

