import tempfile
import unittest
from pathlib import Path

from src.gui.attachment_handler import (
    AttachmentError,
    extract_text_from_file,
    format_attachment_context,
    format_user_text_with_attachments,
    list_supported_files_in_folder,
    validate_attachment_path,
)


class AttachmentHandlerTests(unittest.TestCase):
    def test_extract_text_from_plain_file(self) -> None:
        path = Path(tempfile.mkdtemp()) / "note.txt"
        path.write_text("hello attachment", encoding="utf-8")

        attachment = extract_text_from_file(str(path))

        self.assertEqual(attachment.filename, "note.txt")
        self.assertEqual(attachment.file_type, "txt")
        self.assertEqual(attachment.extracted_text, "hello attachment")

    def test_extract_text_from_pptx_file(self) -> None:
        path = Path(tempfile.mkdtemp()) / "slides.pptx"
        _write_sample_pptx(path, "Quarterly Review", "Revenue increased")

        attachment = extract_text_from_file(str(path))

        self.assertEqual(attachment.filename, "slides.pptx")
        self.assertEqual(attachment.file_type, "pptx")
        self.assertIn("[Slide 1]", attachment.extracted_text)
        self.assertIn("Quarterly Review", attachment.extracted_text)
        self.assertIn("Revenue increased", attachment.extracted_text)

    def test_rejects_unsupported_file_type(self) -> None:
        path = Path(tempfile.mkdtemp()) / "archive.bin"
        path.write_bytes(b"data")

        with self.assertRaises(AttachmentError):
            extract_text_from_file(str(path))

    def test_validate_attachment_path_does_not_extract_content(self) -> None:
        path = Path(tempfile.mkdtemp()) / "note.txt"
        path.write_text("hello attachment", encoding="utf-8")

        self.assertEqual(validate_attachment_path(str(path)), path.resolve())

    def test_list_supported_files_in_folder_finds_only_selected_folder_files(self) -> None:
        root = Path(tempfile.mkdtemp())
        note = root / "note.txt"
        note.write_text("hello", encoding="utf-8")
        slides = root / "slides.pptx"
        slides.write_bytes(b"placeholder")
        ignored = root / "archive.bin"
        ignored.write_bytes(b"data")
        nested = root / "docs"
        nested.mkdir()
        nested_note = nested / "nested.md"
        nested_note.write_text("skip me", encoding="utf-8")

        self.assertEqual(list_supported_files_in_folder(str(root)), [note.resolve(), slides.resolve()])

    def test_format_user_text_with_attachments(self) -> None:
        context = format_attachment_context(
            [
                {
                    "filename": "note.txt",
                    "file_type": "txt",
                    "extracted_text": "attachment text",
                }
            ]
        )

        prompt = format_user_text_with_attachments("summarize", context)

        self.assertIn("File name: note.txt", prompt)
        self.assertIn("attachment text", prompt)
        self.assertTrue(prompt.endswith("User message:\nsummarize"))

def _write_sample_pptx(path: Path, title: str, body: str) -> None:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise unittest.SkipTest("python-pptx is not installed") from exc

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body
    presentation.save(str(path))


if __name__ == "__main__":
    unittest.main()
