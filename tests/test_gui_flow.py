import os
import tempfile
import time
import unittest
import yaml
from pathlib import Path
from threading import Lock
from unittest.mock import patch

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")

from PySide6.QtWidgets import QApplication

from src.gui.config import AppConfig
from src.gui.console_session import ConsoleConfig, ConsoleSessionError
from src.gui.database import ChatRepository
from src.gui.gui import MainWindow
from src.gui.llm_client import OpenAICompatibleClient
from src.gui.token_handler import ModelPrompt


class FakeConsole:
    def __init__(self) -> None:
        self.config = ConsoleConfig("/bin/echo", "/tmp/model.gguf", ctx_size=2048)
        self.prompts: list[str | ModelPrompt] = []
        self.stopped = False
        self._lock = Lock()

    def start(self) -> None:
        self.stopped = False

    def stop(self) -> None:
        self.stopped = True

    def stop_generation(self) -> None:
        self.stopped = True

    def ask(self, prompt: str | ModelPrompt) -> str:
        with self._lock:
            self.prompts.append(prompt)
        for _ in range(20):
            if self.stopped:
                self.stopped = False
                raise ConsoleSessionError("Generation stopped.")
            time.sleep(0.005)
        if isinstance(prompt, ModelPrompt):
            user_text = str(prompt.messages[-1]["content"])
        else:
            user_text = str(prompt)
        return f"assistant answer: {user_text}"


def process_events(app: QApplication, cycles: int = 80) -> None:
    for _ in range(cycles):
        app.processEvents()
        time.sleep(0.005)


def process_until_idle(app: QApplication, window: MainWindow, cycles: int = 120) -> None:
    for _ in range(cycles):
        app.processEvents()
        if not window._active_generations:
            return
        time.sleep(0.005)


class GuiFlowTests(unittest.TestCase):
    def setUp(self) -> None:
        self.app = QApplication.instance() or QApplication([])

    def test_stop_then_next_send_works(self) -> None:
        console = FakeConsole()
        db_path = Path(tempfile.mkdtemp()) / "app.db"
        window = MainWindow(
            console,
            ChatRepository(str(db_path)),
            AppConfig(llama_cli_path="/bin/echo", model_path="/tmp/model.gguf", backend="cli"),
        )

        window.input_edit.setPlainText("first")
        window.on_send()
        process_events(self.app, 5)
        window.on_stop_generation()
        process_events(self.app)

        self.assertIn("Generation stopped.", window.chat_view.toPlainText())
        self.assertTrue(window.send_btn.isEnabled())

        window.input_edit.setPlainText("second")
        window.on_send()
        process_events(self.app)

        self.assertIn("assistant answer", window.chat_view.toPlainText())
        prompt = console.prompts[-1]
        self.assertIsInstance(prompt, ModelPrompt)
        assert isinstance(prompt, ModelPrompt)
        self.assertEqual(prompt.messages[-1]["role"], "user")
        self.assertIn("second", str(prompt.messages[-1]["content"]))
        self.assertIn("second<end_of_turn>", prompt.completion_prompt)
        self.assertTrue(prompt.completion_prompt.endswith("<start_of_turn>model"))

        window.close()
        process_events(self.app, 5)

    def test_generation_finishes_in_original_session_after_switch(self) -> None:
        console = FakeConsole()
        db_path = Path(tempfile.mkdtemp()) / "app.db"
        repository = ChatRepository(str(db_path))
        first_session_id = repository.create_session("First")
        second_session_id = repository.create_session("Second")
        window = MainWindow(
            console,
            repository,
            AppConfig(llama_cli_path="/bin/echo", model_path="/tmp/model.gguf", backend="cli"),
        )

        window.current_session_id = first_session_id
        window._load_session_messages(first_session_id)
        window.input_edit.setPlainText("first prompt")
        window.on_send()
        process_events(self.app, 5)

        window.current_session_id = second_session_id
        window._load_session_messages(second_session_id)
        process_until_idle(self.app, window)

        first_messages = repository.get_messages(first_session_id)
        second_messages = repository.get_messages(second_session_id)
        self.assertEqual(first_messages[-1]["role"], "assistant")
        self.assertEqual(first_messages[-1]["content"], "assistant answer: first prompt")
        self.assertFalse(any(message["role"] == "assistant" for message in second_messages))
        self.assertNotIn("assistant answer: first prompt", window.chat_view.toPlainText())

        window.close()
        process_events(self.app, 5)

    def test_two_sessions_can_generate_at_once(self) -> None:
        console = FakeConsole()
        db_path = Path(tempfile.mkdtemp()) / "app.db"
        repository = ChatRepository(str(db_path))
        first_session_id = repository.create_session("First")
        second_session_id = repository.create_session("Second")
        window = MainWindow(
            console,
            repository,
            AppConfig(llama_cli_path="/bin/echo", model_path="/tmp/model.gguf", backend="cli"),
        )

        window.current_session_id = first_session_id
        window._load_session_messages(first_session_id)
        window.input_edit.setPlainText("first prompt")
        window.on_send()
        process_events(self.app, 5)

        window.current_session_id = second_session_id
        window._load_session_messages(second_session_id)
        window.input_edit.setPlainText("second prompt")
        window.on_send()
        process_until_idle(self.app, window)

        first_messages = repository.get_messages(first_session_id)
        second_messages = repository.get_messages(second_session_id)
        self.assertEqual(first_messages[-1]["content"], "assistant answer: first prompt")
        self.assertEqual(second_messages[-1]["content"], "assistant answer: second prompt")
        self.assertTrue(window.send_btn.isEnabled())

        window.close()
        process_events(self.app, 5)

    def test_send_uses_recent_message_window(self) -> None:
        console = FakeConsole()
        db_path = Path(tempfile.mkdtemp()) / "app.db"
        repository = ChatRepository(str(db_path))
        session_id = repository.create_session("Long chat")
        for index in range(6):
            repository.add_message(session_id, "user", f"user {index}")
            repository.add_message(session_id, "assistant", f"assistant {index}")

        window = MainWindow(
            console,
            repository,
            AppConfig(
                llama_cli_path="/bin/echo",
                model_path="/tmp/model.gguf",
                backend="cli",
                recent_message_limit=3,
            ),
        )
        window.current_session_id = session_id
        window.input_edit.setPlainText("current")

        window.on_send()
        process_events(self.app)

        prompt = console.prompts[-1]
        self.assertIsInstance(prompt, ModelPrompt)
        assert isinstance(prompt, ModelPrompt)
        prompt_text = "\n".join(message["content"] for message in prompt.messages)
        self.assertNotIn("user 0", prompt_text)
        self.assertNotIn("assistant 0", prompt_text)
        self.assertIn("assistant 5", prompt_text)
        self.assertIn("current", prompt_text)
        self.assertIn("Prompt context was shortened", window.chat_view.toPlainText())

        window.close()
        process_events(self.app, 5)

    def test_attach_folder_replaces_existing_attachment_list(self) -> None:
        console = FakeConsole()
        db_path = Path(tempfile.mkdtemp()) / "app.db"
        first_folder = Path(tempfile.mkdtemp())
        second_folder = Path(tempfile.mkdtemp())
        first_file = first_folder / "first.txt"
        second_file = second_folder / "second.pptx"
        first_file.write_text("first", encoding="utf-8")
        _write_sample_pptx(second_file)

        window = MainWindow(
            console,
            ChatRepository(str(db_path)),
            AppConfig(llama_cli_path="/bin/echo", model_path="/tmp/model.gguf", backend="cli"),
        )

        with patch("src.gui.gui.QFileDialog.getExistingDirectory", return_value=str(first_folder)):
            window.on_attach_files()
        self.assertEqual(window._attached_file_paths, [str(first_file.resolve())])

        with patch("src.gui.gui.QFileDialog.getExistingDirectory", return_value=str(second_folder)):
            window.on_attach_files()

        self.assertEqual(window._attached_file_paths, [str(second_file.resolve())])
        self.assertEqual(window.attachment_list.count(), 1)
        self.assertIn("second.pptx", window.attachment_list.item(0).text())

        window.close()
        process_events(self.app, 5)

    def test_send_includes_retrieved_context_from_attached_files(self) -> None:
        console = FakeConsole()
        db_path = Path(tempfile.mkdtemp()) / "app.db"
        attached_folder = Path(tempfile.mkdtemp())
        attached_file = attached_folder / "notes.txt"
        attached_file.write_text("alpha semantic match\nbeta context", encoding="utf-8")

        window = MainWindow(
            console,
            ChatRepository(str(db_path)),
            AppConfig(
                llama_cli_path="/bin/echo",
                model_path="/tmp/model.gguf",
                backend="cli",
                openai_base_url="http://example.test/v1",
                openai_model="embedding-model",
                rag_top_k=2,
                rag_min_score=0.1,
            ),
        )
        window._attached_file_paths = [str(attached_file.resolve())]
        window._attachment_folder_roots[str(attached_file.resolve())] = str(attached_folder.resolve())
        window._refresh_attachment_list()

        def fake_embeddings(self, inputs, model=None):
            vectors = []
            for text in inputs:
                lowered = str(text).lower()
                if "question about alpha" in lowered or "alpha semantic match" in lowered:
                    vectors.append([1.0, 0.0])
                else:
                    vectors.append([0.0, 1.0])
            return vectors

        with patch.object(OpenAICompatibleClient, "embeddings", new=fake_embeddings):
            window.input_edit.setPlainText("question about alpha")
            window.on_send()
            process_events(self.app)

        prompt = console.prompts[-1]
        self.assertIsInstance(prompt, ModelPrompt)
        assert isinstance(prompt, ModelPrompt)
        system_messages = [message["content"] for message in prompt.messages if message["role"] == "system"]
        self.assertTrue(any("Relevant retrieved context:" in message for message in system_messages))
        self.assertTrue(any("notes.txt" in message for message in system_messages))
        self.assertTrue(any("alpha semantic match" in message for message in system_messages))

        window.close()
        process_events(self.app, 5)

    def test_last_working_folder_is_restored_on_startup(self) -> None:
        console = FakeConsole()
        db_path = Path(tempfile.mkdtemp()) / "app.db"
        config_path = Path(tempfile.mkdtemp()) / "config.yaml"
        attached_folder = Path(tempfile.mkdtemp())
        attached_file = attached_folder / "notes.txt"
        attached_file.write_text("hello", encoding="utf-8")
        config_path.write_text("backend: cli\nlast_working_folder: ''\n", encoding="utf-8")

        window = MainWindow(
            console,
            ChatRepository(str(db_path)),
            AppConfig(
                llama_cli_path="/bin/echo",
                model_path="/tmp/model.gguf",
                backend="cli",
                config_path=str(config_path),
                last_working_folder=str(attached_folder.resolve()),
            ),
        )

        self.assertEqual(window._active_attachment_folder(), str(attached_folder.resolve()))
        self.assertEqual(window.attachment_list.count(), 1)
        self.assertIn("notes.txt", window.attachment_list.item(0).text())

        saved = yaml.safe_load(config_path.read_text(encoding="utf-8"))
        self.assertEqual(saved["last_working_folder"], str(attached_folder.resolve()))

        window.close()
        process_events(self.app, 5)

    def test_missing_last_working_folder_is_cleared(self) -> None:
        console = FakeConsole()
        db_path = Path(tempfile.mkdtemp()) / "app.db"
        config_path = Path(tempfile.mkdtemp()) / "config.yaml"
        missing_folder = Path(tempfile.mkdtemp()) / "missing"
        config_path.write_text("backend: cli\nlast_working_folder: ''\n", encoding="utf-8")

        window = MainWindow(
            console,
            ChatRepository(str(db_path)),
            AppConfig(
                llama_cli_path="/bin/echo",
                model_path="/tmp/model.gguf",
                backend="cli",
                config_path=str(config_path),
                last_working_folder=str(missing_folder),
            ),
        )

        self.assertIsNone(window._active_attachment_folder())
        saved = yaml.safe_load(config_path.read_text(encoding="utf-8"))
        self.assertEqual(saved["last_working_folder"], "")

        window.close()
        process_events(self.app, 5)

def _write_sample_pptx(path: Path) -> None:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise unittest.SkipTest("python-pptx is not installed") from exc

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Deck"
    slide.placeholders[1].text = "Body"
    presentation.save(str(path))


if __name__ == "__main__":
    unittest.main()
